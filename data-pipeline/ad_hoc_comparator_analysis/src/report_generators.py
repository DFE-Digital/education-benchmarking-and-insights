import os
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.text import PP_ALIGN
from io import BytesIO
import matplotlib.pyplot as plt
from src.functions import plot_comparator_map
from src.functions import comp_barplot
from src.functions import overall_summary_dataframe


#add charts to a powerpoint deck
def add_chart_to_pptx(slide, fig, left, top, width):
    img_stream = BytesIO()
    if hasattr(fig, 'fig'): fig.fig.savefig(img_stream, format='png', bbox_inches='tight', dpi=50)
    else: fig.savefig(img_stream, format='png', bbox_inches='tight', dpi=50)
    slide.shapes.add_picture(img_stream, Cm(left), Cm(top), width=Cm(width))
    plt.close()



def add_table_to_pptx(slide, df, left_cm, top_cm, width_cm):
    # Takes a dataframe and puts it into a powerpoint slide as a table
    rows, cols = df.shape
    table_shape = slide.shapes.add_table(rows + 1, cols, Cm(left_cm), Cm(top_cm), Cm(width_cm), Cm(0.8))
    table = table_shape.table

    # Header Row
    for col_idx, col_name in enumerate(df.columns):
        cell = table.cell(0, col_idx)
        cell.text = str(col_name)
        _apply_styling(cell, is_header=True)

    # Data Rows
    for row_idx in range(rows):
        for col_idx in range(cols):
            cell = table.cell(row_idx + 1, col_idx)
            val = df.iloc[row_idx, col_idx]
            
            # Formatting numbers
            cell.text = f"{val:,.2f}" if isinstance(val, (float, int)) else str(val)
            
            # Right-align numbers, Left-align text
            alignment = PP_ALIGN.RIGHT if isinstance(val, (int, float)) else PP_ALIGN.LEFT
            _apply_styling(cell, is_header=False, align=alignment)

    return table_shape

def _apply_styling(cell, is_header=False, align=PP_ALIGN.CENTER):
    paragraph = cell.text_frame.paragraphs[0]
    paragraph.alignment = align
    font = paragraph.font
    font.name = 'Arial'
    font.size = Pt(12) if not is_header else Pt(14)
    font.bold = is_header


#This produces the actual powerpoint deck, using the functions above
def generate_pptx(target_urn, school_name, processed_results, flat_results_df, output_path):
    prs = Presentation()
    # Change aspect ratio from old school 4:3 to standard widescreen 16:9
    prs.slide_width = Cm(33.876)
    prs.slide_height = Cm(19.05)
    # Simplified title and content loops

    info_slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    info_title_box1 = info_slide1.shapes.add_textbox(Cm(1.27), Cm(1), Cm(22), Cm(2.6))
    info_title_box1.text_frame.text = f"URN: {target_urn} {school_name}"

    info_content_box1 = info_slide1.shapes.add_textbox(Cm(1.27), Cm(2), Cm(30), Cm(18))
    info_content_box1.text_frame.word_wrap = True
    info_content_box1.text_frame.paragraphs[0].font.size = Pt(9)
    info_content_box1.text_frame.text = "\nThis slidedeck compares four different algorithms for comparator group selection." \
    "\n"\
    "\nFBIT calculates two comparator groups for each school: Pupil, and Building." \
    "\n"\
    "\nIn the existing algorithm Pupils is based on Number of Pupils, Percentage SEN, and Percentage Free School Meals, and Buildings is based on Total Internal Floor Area and Average Age Score. A series of filters are applied based on phase, finance type, geography etc. before using those features to calculate a similarity measure."\
    "\n"\
    "\nThe existing algorithm weights features differently in that calculation, e.g. pupil numbers is responsible for 50% of the similarity measure."\
    "\n"\
    
    info_slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    info_title_box2 = info_slide2.shapes.add_textbox(Cm(1.27), Cm(1), Cm(22), Cm(2.6))
    info_title_box2.text_frame.text = f"URN: {target_urn} {school_name}"
    
    info_content_box2 = info_slide2.shapes.add_textbox(Cm(1.27), Cm(2), Cm(30), Cm(18))
    info_content_box2.text_frame.word_wrap = True
    info_content_box2.text_frame.paragraphs[0].font.size = Pt(9)
    info_content_box2.text_frame.text = "\nThe algorithms compared are defined as follows:"\
    "\n     - 1 baseline: the method currently used by FBIT" \
    "\n     - 2 baseline_geog: same features and weightings as baseline, but uses as the crow flies geographic distance" \
    "\n     - 3 baseline_equal_geog: same as baseline_geog, except the features are equally weighted" \
    "\n     - 4 additional_features_geog: all features shown are used in the comparator group selection,"\
    "\n         along with geographic distance, and are equally weighted"\
    "\n"\
    "\nThe full list of additional features are shown in the tables on the next slide, and in the barplots on each of the subsequent slides"\
    "\n"\
    "\nAfter the summary slide, there are two slides for each algorithm choice - one for the Pupil comparator group, and one for the Building group"\
    "\n"\
    "\nEach contains a map showing the geographical distribution of the comparators (URNs shown on map, target school in red) and boxplots showing the distribution of values of the features within the group"\
    "\n - target school in orange, red dotted lines showing threshold values used to include comparator schools in the FBIT rag calculation."

    summary_slide = prs.slides.add_slide(prs.slide_layouts[6])
    summary_title_box = summary_slide.shapes.add_textbox(Cm(1.27), Cm(1), Cm(22), Cm(2.6))
    summary_title_box.text_frame.text = f"URN: {target_urn} {school_name}"

    pupil_summary_dataframe = overall_summary_dataframe(flat_results_df, target_urn, comp_type='pupil', 
                                                        features=['distance_km', 'Number of pupils', 'Percentage SEN', 
                                                                  'Percentage Free school meals', 
                                                                  'unfilled_places_count', 'pupils_over_capacity_count',
                                                                  'RuralScore', 'sparse'])
    
    building_summary_dataframe = overall_summary_dataframe(flat_results_df, target_urn, comp_type='building', 
                                                           features=['distance_km', 'Total Internal Floor Area', 
                                                                     'OldestBuildingAge', 'NewestBuildingAge', 
                                                                     'Age Average Score', 'BuildingCount', 'SplitSiteScore'])

    pupil_summary_textbox = summary_slide.shapes.add_textbox(Cm(1.27), Cm(2), Cm(22), Cm(0.8))
    pupil_summary_textbox.text_frame.text = "Pupil Comparators (median by algorithm, and actual/target)"

    add_table_to_pptx(summary_slide, pupil_summary_dataframe, left_cm=1.27, top_cm=3.5, width_cm=30)

    building_summary_textbox = summary_slide.shapes.add_textbox(Cm(1.27), Cm(11), Cm(22), Cm(0.8))
    building_summary_textbox.text_frame.text = "Building Comparators (median by algorithm, and actual/target)"
    add_table_to_pptx(summary_slide, building_summary_dataframe, left_cm=1.27, top_cm=12, width_cm=30)

    for algo_name, dfs in processed_results.items():
        
        for comparator_group, dataset in dfs.items():

            slide = prs.slides.add_slide(prs.slide_layouts[6])

            title_box = slide.shapes.add_textbox(Cm(1.27), Cm(1), Cm(20), Cm(2.6))
            title_box.text_frame.text = f"Algorithm {algo_name} | Comparator Group {comparator_group}  | URN: {target_urn} | {dataset[dataset["URN"] == target_urn]["SchoolName"].iloc[0]} \n | {dataset[dataset["URN"] == target_urn]["SchoolPhaseType"].iloc[0]} | {dataset[dataset["URN"] == target_urn]["FinanceType"].iloc[0]} | PFI {dataset[dataset["URN"] == target_urn]["Is PFI"].iloc[0]} | {dataset[dataset["URN"] == target_urn]["Boarders (name)"].iloc[0]}"
            # Logic to add map and bar charts using add_chart_to_pptx

            #geospatial map
            fig_map = plot_comparator_map(dataset,target_urn)
            add_chart_to_pptx(slide, fig_map, left=1.3, top=3.7, width=11)
            plt.close(fig_map)

            #feature barcharts
            if comparator_group == "pupil":
                fig_pupils = comp_barplot(dataset, target_urn, "Number of pupils", threshold=0.25) #RAG threshold for pupil close comparators is +/- 25%
                add_chart_to_pptx(slide, fig_pupils, left=14.3, top=3.9, width=5.2)
                
                fig_sen = comp_barplot(dataset, target_urn, "Percentage SEN", threshold=0.1) #RAG threshold for pupil close comparators is +/- 10%
                add_chart_to_pptx(slide, fig_sen, left=20.8, top=3.9, width=5.2)
                
                fig_fsm = comp_barplot(dataset, target_urn, "Percentage Free school meals", threshold=0.05) #RAG threshold for pupil close comparators is +/- 5%
                add_chart_to_pptx(slide, fig_fsm, left=27.3, top=3.9, width=5.2)
                
                fig_undercapacity = comp_barplot(dataset, target_urn, "unfilled_places_count", threshold='none')
                add_chart_to_pptx(slide, fig_undercapacity, left=14.3, top=9.1, width=5.2)
                
                fig_overcapacity = comp_barplot(dataset, target_urn, "pupils_over_capacity_count", threshold='none')
                add_chart_to_pptx(slide, fig_overcapacity, left=20.8, top=9.1, width=5.2)
                
                fig_rurality = comp_barplot(dataset, target_urn, "RuralScore", threshold='none')
                add_chart_to_pptx(slide, fig_rurality, left=27.3, top=9.1, width=5.2)

                fig_sparsity = comp_barplot(dataset, target_urn, "sparse", threshold='none')
                add_chart_to_pptx(slide, fig_sparsity, left=14.3, top=14.3, width=5.2)

                fig_distance = comp_barplot(dataset, target_urn, "distance_km", threshold='none')
                add_chart_to_pptx(slide, fig_distance, left=27.3, top=14.3, width=5.2)
                

            elif comparator_group == "building":
                fig_gifa = comp_barplot(dataset, target_urn, "Total Internal Floor Area", threshold=0.1) #RAG threshold for building close comparators is +/- 10%
                add_chart_to_pptx(slide, fig_gifa, left=14.3, top=3.9, width=5.2)
                
                fig_old = comp_barplot(dataset, target_urn, "OldestBuildingAge", threshold='none')
                add_chart_to_pptx(slide, fig_old, left=20.8, top=3.9, width=5.2)
                
                fig_new = comp_barplot(dataset, target_urn, "NewestBuildingAge", threshold='none')
                add_chart_to_pptx(slide, fig_new, left=27.3, top=3.9, width=5.2)
                
                fig_age = comp_barplot(dataset, target_urn, "Age Average Score", threshold=0.2) #RAG threshold for building close comparators is +/- 20%
                add_chart_to_pptx(slide, fig_age, left=14.3, top=9.1, width=5.2)
                
                fig_buildings = comp_barplot(dataset, target_urn, "BuildingCount", threshold='none')
                add_chart_to_pptx(slide, fig_buildings, left=20.8, top=9.1, width=5.2)

                fig_splitsite = comp_barplot(dataset, target_urn, "SplitSiteScore", threshold='none')
                add_chart_to_pptx(slide, fig_splitsite, left=27.3, top=9.1, width=5.2)

                fig_distance = comp_barplot(dataset, target_urn, "distance_km", threshold='none')
                add_chart_to_pptx(slide, fig_distance, left=14.3, top=14.3, width=5.2)
                
            pass 
        prs.save(output_path)